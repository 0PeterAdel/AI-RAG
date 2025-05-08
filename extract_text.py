import os
import json
from pathlib import Path
import pdfplumber
import docx
from pptx import Presentation
import pandas as pd

# مسار المجلد الذي يحوي كل الملفات
ROOT_DIR = Path("Data")

# مجلد حفظ النصوص
OUTPUT_DIR = Path("02.Raw-Text")
OUTPUT_DIR.mkdir(exist_ok=True)

# دوال الاستخلاص
def extract_from_pdf(path):
    text = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            text.append(page.extract_text() or "")
    return "\n".join(text)

def extract_from_docx(path):
    doc = docx.Document(path)
    return "\n".join(p.text for p in doc.paragraphs)

def extract_from_pptx(path):
    prs = Presentation(path)
    slides = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slides.append(shape.text)
    return "\n".join(slides)

def extract_from_csv(path):
    df = pd.read_csv(path, dtype=str, encoding_errors="ignore")
    # ندمج كل صف في سطر نصي واحد بفواصل
    return "\n".join(df.astype(str).agg(" | ".join, axis=1))

def extract_from_excel(path):
    df = pd.read_excel(path, dtype=str)
    return "\n".join(df.astype(str).agg(" | ".join, axis=1))

# خريطة الامتدادات إلى الدوال
HANDLERS = {
    ".pdf": extract_from_pdf,
    ".docx": extract_from_docx,
    ".pptx": extract_from_pptx,
    ".csv": extract_from_csv,
    ".xlsx": extract_from_excel,
}

# مرور على كل الملفات
for subdir, _, files in os.walk(ROOT_DIR):
    # نتجاهل مجلد الإخراج نفسه
    if OUTPUT_DIR.name in subdir:
        continue

    for fname in files:
        path = Path(subdir) / fname
        ext = path.suffix.lower()
        if ext in HANDLERS:
            print(f"Extracting {path}…")
            try:
                text = HANDLERS[ext](path)
                # نحفظ لكل ملف JSON باسم مطابق
                out_path = OUTPUT_DIR / (path.relative_to(ROOT_DIR).with_suffix(".json").as_posix())
                out_path.parent.mkdir(parents=True, exist_ok=True)
                json.dump({
                    "source": str(path),
                    "text": text
                }, open(out_path, "w", encoding="utf-8"), ensure_ascii=False, indent=2)
            except Exception as e:
                print(f" ⚠️ Failed {path}: {e}")
